from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor

from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_colored_background(slide, color):
    """Add a colored background to a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_and_content(slide, title_text, content_text, color):
    """Add a title and content to a slide with color customization."""
    add_colored_background(slide, color)

    # Add title
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Add content
    content_box = slide.placeholders[1]
    content_box.text = content_text
    for paragraph in content_box.text_frame.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.alignment = PP_ALIGN.LEFT

def create_colorful_presentation():
    # Create presentation object
    prs = Presentation()

    # ------------------------ Title Slide ------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    add_colored_background(slide, RGBColor(0, 102, 204))
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "SonarQube: An Overview of Code Quality and Security"
    subtitle.text = "Enhancing Software Quality Through Continuous Inspection"

    # ------------------------ Slide 2: Agenda ------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_content(slide, "Agenda", """\
1. Introduction to SonarQube
2. Why Code Quality Matters
3. Key Features of SonarQube
4. How SonarQube Works
5. Benefits of Using SonarQube
6. Use Cases & Integrations
7. Demo or Example Workflow
8. Conclusion & Q&A""", RGBColor(102, 51, 153))

    # ------------------------ Slide 3: Introduction to SonarQube ------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_content(slide, "What is SonarQube?", """\
- An open-source platform for continuous inspection of code quality and security.
- Detects code smells, bugs, and security vulnerabilities.
- Supports 30+ programming languages (Java, C#, Python, etc.).""", RGBColor(0, 153, 153))

    # ------------------------ Slide 4: Why Code Quality Matters ------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_content(slide, "Why Focus on Code Quality?", """\
- Reduces technical debt and maintenance costs.
- Improves readability, reliability, and security.
- Enhances team productivity and software stability.
- \"80% of maintenance costs are spent on poorly written code.\"""", RGBColor(204, 51, 51))

    # ------------------------ Slide 5: Key Features of SonarQube ------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_content(slide, "Key Features", """\
- Static Code Analysis for bugs and vulnerabilities.
- Code Smell Detection for maintainability issues.
- Security Hotspot Detection for potential security risks.
- Quality Gates to ensure code meets standards before release.
- Reports & Dashboards for comprehensive insights.""", RGBColor(51, 102, 204))

    # ------------------------ Slide 6: How SonarQube Works ------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_content(slide, "How SonarQube Works", """\
1. Code Scanning: SonarQube scans your code for issues.
2. Analysis Engine: Rulesets analyze for quality and security flaws.
3. Reports & Feedback: Generates detailed reports.
4. Continuous Integration: Works with CI/CD tools (Jenkins, GitHub Actions).""", RGBColor(0, 102, 51))

    # ------------------------ Slide 7: Benefits of SonarQube ------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_content(slide, "Benefits of SonarQube", """\
- Ensures continuous code quality.
- Reduces technical debt.
- Promotes best coding practices.
- Integrates with popular CI/CD tools.
- Supports team collaboration and accountability.""", RGBColor(153, 51, 102))

    # ------------------------ Slide 8: Conclusion ------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_content(slide, "Conclusion", """\
- SonarQube helps you write cleaner, safer, and more maintainable code.
- Next Steps: Explore SonarQube, integrate it into your workflow, and improve code quality.

Questions? Let’s Discuss!""", RGBColor(51, 51, 153))

    # Save the presentation
    prs.save("colorful_sonarqube_presentation.pptx")

if __name__ == "__main__":
    create_colorful_presentation()
