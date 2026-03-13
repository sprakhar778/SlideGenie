import subprocess
import os
import logging
from src.api.helpers import load_presentation, PRESENTATIONS_DIR

logger = logging.getLogger(__name__)

def generate_presentation_pptx_from_pdf(presentation_id: str) -> str:
    """
    First generates a PDF, then converts it to PPTX using external tools.
    Requires: pip install pdf2image and have LibreOffice installed
    """
    # First generate PDF using your existing PDF function
    from src.helpers.html_to_pdf import generate_presentation_pdf
    pdf_path = generate_presentation_pdf(presentation_id)
    
    output_dir = PRESENTATIONS_DIR
    pptx_path = os.path.join(output_dir, f"{presentation_id}.pptx")
    
    try:
        # Method 1: Using LibreOffice (if installed)
        subprocess.run([
            'libreoffice', '--headless', '--convert-to', 'pptx',
            '--outdir', output_dir, pdf_path
        ], check=True)
        
        # Method 2: Using pdf2image and python-pptx (fallback)
        if not os.path.exists(pptx_path):
            from pdf2image import convert_from_path
            from pptx import Presentation
            from pptx.util import Inches
            import tempfile
            
            images = convert_from_path(pdf_path)
            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            
            for image in images:
                # Save image temporarily
                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_img:
                    image.save(tmp_img.name, 'PNG')
                    
                    # Add slide with image
                    slide_layout = prs.slide_layouts[6]
                    slide = prs.slides.add_slide(slide_layout)
                    slide.shapes.add_picture(
                        tmp_img.name, 
                        Inches(0), 
                        Inches(0), 
                        width=prs.slide_width,
                        height=prs.slide_height
                    )
                    
                    os.unlink(tmp_img.name)
            
            prs.save(pptx_path)
            print("Done",pptx_path)
            
        return os.path.abspath(pptx_path)
        
    except Exception as e:
        logger.error(f"Failed to convert PDF to PPTX: {e}")
        raise


generate_presentation_pptx_from_pdf(presentation_id="ae5be2ba-0cbe-4320-adf4-8d8766102b3a")